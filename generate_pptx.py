import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to 16:9 Widescreen (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6] # Blank layout

    # Dark Theme Colors
    BG_COLOR = RGBColor(11, 15, 25)         # #0B0F19 Dark Navy/Slate
    HEADER_CARD_BG = RGBColor(18, 24, 38)  # #121826 Panel BG
    CYAN_ACCENT = RGBColor(0, 229, 255)     # #00E5FF Bright Cyan
    BLUE_ACCENT = RGBColor(56, 189, 248)    # #38BDF8 Light Blue
    TEXT_WHITE = RGBColor(255, 255, 255)    # White
    TEXT_BODY = RGBColor(226, 232, 240)     # Light Slate
    TEXT_MUTED = RGBColor(148, 163, 184)    # Muted Slate
    TEXT_DIM = RGBColor(100, 116, 139)      # Dim Slate
    AMBER_ACCENT = RGBColor(255, 183, 3)    # Amber

    slides_data = [
        {
            "tag": "SLIDE 01 / 06 | INCIDENT OVERVIEW & CLASSIFICATION",
            "title": "What is NotPetya?",
            "bullets": [
                [("Initial Outbreak: ", True), ("NotPetya was a major cyberattack that began on ", False), ("June 27, 2017", True), (".", False)],
                [("Initial Appearance: ", True), ("Initially appeared to be ransomware, closely mimicking Petya ransomware.", False)],
                [("Ransom Demand: ", True), ("Displayed a red lock-screen requiring a ", False), ("$300 Bitcoin ransom", True), (" payment.", False)],
                [("True Technical Classification: ", True), ("Identified by digital forensics as a ", False), ("destructive wiper", True), (", not genuine ransomware.", False)],
                [("Primary Purpose: ", True), ("Designed strictly for ", False), ("widespread disruption and permanent destruction", True), (" rather than financial profit.", False)],
            ],
            "notes": "Introduce the attack date (June 27, 2017) and explain the key distinction: although it looked like ransomware demanding $300, forensic evidence showed the encryption key was erased immediately, proving its real goal was permanent destruction."
        },
        {
            "tag": "SLIDE 02 / 06 | INITIAL ACCESS & SUPPLY CHAIN VECTOR",
            "title": "Where Did It Start and How?",
            "bullets": [
                [("Ground Zero: ", True), ("The attack primarily originated in ", False), ("Ukraine", True), (".", False)],
                [("Major Ukrainian Targets: ", True), ("Paralyzed government organizations, banks, energy companies, transportation, and media outlets.", False)],
                [("Compromised Infrastructure: ", True), ("Attackers compromised the update servers of ", False), ("M.E.Doc", True), (", a widely used Ukrainian accounting software.", False)],
                [("Malicious Software Update: ", True), ("A backdoor-infected update was automatically distributed to legitimate customers.", False)],
                [("Forensic Significance: ", True), ("Serves as a major real-world example of a catastrophic ", False), ("software supply-chain attack", True), (".", False)],
            ],
            "notes": "M.E.Doc was mandatory for business tax filings in Ukraine. By breaching the update server, attackers weaponized a trusted software channel to infect thousands of enterprise networks automatically."
        },
        {
            "tag": "SLIDE 03 / 06 | PROPAGATION & LATERAL MOVEMENT",
            "title": "How Did NotPetya Spread?",
            "bullets": [
                [("EternalBlue: ", True), ("Exploited CVE-2017-0144 vulnerability in Windows SMBv1 protocol for remote code execution.", False)],
                [("EternalRomance: ", True), ("Utilized another SMBv1 exploitation technique for remote privilege escalation.", False)],
                [("Credential Theft: ", True), ("Scraped active administrator credentials and password hashes directly from infected system memory (LSASS).", False)],
                [("PsExec: ", True), ("Used legitimate Windows administration utility (PsExec) to execute payloads remotely on connected hosts.", False)],
                [("WMI (Windows Management Instrumentation): ", True), ("Used native WMI commands for stealthy remote execution.", False)],
                [("Rapid Propagation: ", True), ("Once inside an organization, NotPetya could spread rapidly across connected domain systems within minutes.", False)],
            ],
            "notes": "Emphasize that NotPetya didn't rely on just one exploit. Even on patched systems, it stole domain credentials from infected RAM and used legitimate Windows admin tools (PsExec and WMI) to log in and infect neighboring machines."
        },
        {
            "tag": "SLIDE 04 / 06 | PAYLOAD EXECUTION & SYSTEM DESTRUCTION",
            "title": "What Happened to Infected Systems?",
            "bullets": [
                [("Payload Execution: ", True), ("The malicious update executed NotPetya on the victim's computer with administrative rights.", False)],
                [("Reconnaissance & Movement: ", True), ("Harvested credentials and moved laterally through the internal network.", False)],
                [("MBR Destruction: ", True), ("Targeted and corrupted the ", False), ("Master Boot Record (MBR)", True), (" and Master File Table (MFT).", False)],
                [("System Inoperability: ", True), ("Forced a reboot, displaying a fake chkdsk screen before rendering the operating system completely unbootable.", False)],
                [("Ransom Message: ", True), ("Displayed a ransom message demanding $300 in Bitcoin for file recovery.", False)],
                [("Irrecoverable Recovery Mechanism: ", True), ("Victims did not have a reliable method to recover files by paying; key was discarded, proving intentional destruction.", False)],
            ],
            "notes": "Walk through the destruction process: NotPetya triggered a reboot, showed a fake CHKDSK disk repair screen while encrypting file tables, and corrupted sector 0. Reverse engineering confirmed the installation key was randomly generated garbage, making recovery impossible."
        },
        {
            "tag": "SLIDE 05 / 06 | GLOBAL IMPACT & DIGITAL FORENSICS INVESTIGATION",
            "title": "Major Impact and Digital Forensic Investigation",
            "bullets": [
                [("Major Organizations Affected:", True)],
                [("   • Maersk: ", True), ("World's largest shipping container firm; IT infrastructure paralyzed (45,000 PCs destroyed).", False)],
                [("   • Merck: ", True), ("Global pharmaceutical leader; vaccine manufacturing disrupted.", False)],
                [("   • FedEx / TNT Express: ", True), ("European parcel delivery network crippled.", False)],
                [("   • Mondelez & Saint-Gobain: ", True), ("Multinational food and industrial manufacturing halted.", False)],
                [("Global Economic Loss: ", True), ("Disruption to shipping and logistics led to estimated global losses of ", False), ("more than $10 billion", True), (".", False)],
                [("Digital Forensic Artifacts Examined: ", True), ("Analyzed malware code/behavior, Windows Event Logs (Event IDs 4688, 7045, 4624), network logs, M.E.Doc updates, credential artifacts, disk images, memory dumps, and MBR modifications.", False)],
            ],
            "notes": "Highlight the scale: over $10 billion in damage across global supply chains. Forensic examiners proved the wiper behavior by inspecting raw disk sectors, memory dumps, Windows event logs, and M.E.Doc updater logs."
        },
        {
            "tag": "SLIDE 06 / 06 | ATTRIBUTION, LESSONS & CONCLUSION",
            "title": "Attribution, Conclusion and Lessons Learned",
            "bullets": [
                [("Attribution: ", True), ("In 2018, US, UK, and Australia publicly attributed NotPetya to Russian military intelligence (", False), ("GRU / Sandworm", True), ("). Russia denied responsibility.", False)],
                [("Key Cybersecurity Lessons: ", True), ("Secure software supply chain, patch SMB vulnerabilities quickly, enforce network segmentation, protect admin credentials, apply least-privilege access, maintain offline backups, monitor unusual network activity (PsExec/WMI), and maintain forensic readiness.", False)],
                [("Final Conclusion: ", True), ("NotPetya demonstrated that a cyberattack can begin through a trusted software update, spread rapidly through an organization, and cause enormous global damage. Although it appeared to be ransomware, forensic evidence showed that its real capability and purpose were largely destructive.", False)],
            ],
            "notes": "Conclude with the main takeaway: NotPetya redefined cyber threat models. It proved that trusted updates can be backdoored and that digital forensics is critical to distinguish fake ransomware from state-sponsored cyber warfare."
        }
    ]

    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Set solid dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

        # Header Box Card (Background panel shape)
        shapes = slide.shapes
        
        # Tag text box (Category)
        txBox_tag = shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.4))
        tf_tag = txBox_tag.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = slide_data["tag"]
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = CYAN_ACCENT
        p_tag.font.name = "Arial"

        # Title text box
        txBox_title = shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.733), Inches(0.8))
        tf_title = txBox_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = slide_data["title"]
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE
        p_title.font.name = "Arial"

        # Separator line text box / border visual representation
        line_box = shapes.add_textbox(Inches(0.8), Inches(1.65), Inches(11.733), Inches(0.1))
        tf_line = line_box.text_frame
        p_line = tf_line.paragraphs[0]
        p_line.text = "─" * 85
        p_line.font.size = Pt(10)
        p_line.font.color.rgb = TEXT_DIM

        # Content bullets text box
        txBox_content = shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0))
        tf_content = txBox_content.text_frame
        tf_content.word_wrap = True
        tf_content.margin_left = Inches(0)
        tf_content.margin_top = Inches(0)

        for i, bullet_chunks in enumerate(slide_data["bullets"]):
            p = tf_content.add_paragraph() if i > 0 else tf_content.paragraphs[0]
            p.space_after = Pt(12)
            p.line_spacing = 1.25

            # Bullet character
            run_bullet = p.add_run()
            run_bullet.text = "▪  "
            run_bullet.font.color.rgb = CYAN_ACCENT
            run_bullet.font.size = Pt(15)
            run_bullet.font.name = "Arial"

            for text, is_bold in bullet_chunks:
                run = p.add_run()
                run.text = text
                run.font.size = Pt(15)
                run.font.name = "Arial"
                if is_bold:
                    run.font.bold = True
                    run.font.color.rgb = BLUE_ACCENT
                else:
                    run.font.bold = False
                    run.font.color.rgb = TEXT_BODY

        # Add Speaker Notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = slide_data["notes"]

    output_path = os.path.join(os.getcwd(), "NotPetya_Digital_Forensics_Presentation.pptx")
    prs.save(output_path)
    print(f"SUCCESS: Created PowerPoint presentation at {output_path}")

if __name__ == "__main__":
    create_presentation()
